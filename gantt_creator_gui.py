"""
gantt_creator_gui.py
---------------------

This module provides a simple graphical interface for creating and
editing Gantt charts based on an Excel schedule template.  It relies
on ``tkinter`` for the GUI and Matplotlib for rendering the chart.  The
interface allows the user to:

* Load an Excel file containing timeline and heat‑map entries.
* View and edit individual task titles and bar colours via a form.
* Adjust the horizontal zoom factor with a slider to compress or
  expand the timeline.
* Generate and display the Gantt chart directly inside the window.
* Increase the size of milestone diamonds by specifying a marker size.
* Save the chart as a PNG file or export it to a PowerPoint slide.

Due to the limitations of this environment, we cannot test the GUI
here.  However, this script is designed to run on a local machine
where ``tkinter`` is available (as verified by importing ``tkinter``).

Author: ChatGPT
Date: 2025‑08‑02
"""

from __future__ import annotations

import os
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from matplotlib.patches import Patch
from matplotlib.lines import Line2D
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
from matplotlib.colors import to_hex, to_rgb
from datetime import datetime, timedelta
from dateutil.relativedelta import relativedelta
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, colorchooser

# -----------------------------------------------------------------------------
# Utility functions copied and adapted from gantt_creator.py
#

def lighten_color(color: str, amount: float = 0.5) -> str:
    """Lighten the given colour by mixing it with white.

    Parameters
    ----------
    color : str
        A matplotlib‑compatible colour string (hex or named).
    amount : float, optional
        Amount by which to lighten the colour.  0.0 returns the
        original colour, 1.0 returns white.

    Returns
    -------
    str
        A hex colour string representing the lightened colour.
    """
    try:
        c = to_rgb(color)
    except ValueError:
        c = (0.5, 0.5, 0.5)
    r, g, b = c
    r = r + (1.0 - r) * amount
    g = g + (1.0 - g) * amount
    b = b + (1.0 - b) * amount
    return to_hex((r, g, b))


def compute_date_range(timeline_df: pd.DataFrame, heat_df: pd.DataFrame) -> tuple[datetime, datetime]:
    """Compute the overall min/max dates across timeline and heat map sections.
    Adds a small margin for aesthetics.
    """
    dates_start: list[datetime] = []
    dates_end: list[datetime] = []
    if not timeline_df.empty:
        if timeline_df['Date From'].notna().any():
            dates_start.append(timeline_df['Date From'].min())
        if timeline_df['Date To'].notna().any():
            dates_end.append(timeline_df['Date To'].max())
        else:
            dates_end.append(timeline_df['Date From'].max())
    if not heat_df.empty and heat_df['Heat Map Dates'].notna().any():
        dates_start.append(heat_df['Heat Map Dates'].min())
        dates_end.append(heat_df['Heat Map Dates'].max())
    if not dates_start or not dates_end:
        raise ValueError("No valid dates found in dataset.")
    start_date = min(dates_start)
    end_date = max(dates_end)
    total_days = (end_date - start_date).days
    margin_days = max(5, int(total_days * 0.02))
    margin = timedelta(days=margin_days)
    return start_date - margin, end_date + margin


def generate_color_maps(unique_symbols: list[str], heat_lines: list[str]) -> tuple[dict[str, str], dict[str, dict[str, str]], dict[str, str]]:
    """Generate colour mappings for timeline bars, markers and heat map rows."""
    bar_color_defaults = {
        'Blue Bar': '#1f77b4',
        'Grey Bar': '#8c8c8c',
        'White Bar': '#ffffff',
        'Red Bar': '#d62728',
    }
    marker_map: dict[str, dict[str, str]] = {
        'Milestone': {'marker': 'D', 'color': 'black'},
        'Black Spot': {'marker': 'o', 'color': 'black'},
        'Red Spot': {'marker': 'o', 'color': 'red'},
    }
    bar_color_map: dict[str, str] = {}
    palette = plt.cm.tab10.colors
    palette_iter = iter(palette)
    for sym in unique_symbols:
        # If the symbol is a valid hex colour code, use it directly.
        if isinstance(sym, str) and sym.startswith('#') and len(sym) in (7, 9):
            bar_color_map[sym] = sym
            continue
        if sym in bar_color_defaults:
            bar_color_map[sym] = bar_color_defaults[sym]
        elif sym in marker_map:
            # markers handled separately
            continue
        else:
            try:
                bar_color_map[sym] = to_hex(next(palette_iter))
            except StopIteration:
                palette_iter = iter(palette)
                bar_color_map[sym] = to_hex(next(palette_iter))
    heat_palette = plt.cm.tab20.colors
    heat_colour_map: dict[str, str] = {}
    heat_iter = iter(heat_palette)
    for line in heat_lines:
        try:
            heat_colour_map[line] = to_hex(next(heat_iter))
        except StopIteration:
            heat_iter = iter(heat_palette)
            heat_colour_map[line] = to_hex(next(heat_iter))
    return bar_color_map, marker_map, heat_colour_map


def create_gantt_figure(
    df: pd.DataFrame,
    title: str,
    zoom_factor: float = 1.0,
    marker_size: int = 70,
    suppress_duplicate_labels: bool = True,
    label_offset_scale: float = 0.25,
) -> plt.Figure:
    """Create a Matplotlib figure containing the Gantt chart.

    Parameters
    ----------
    df : pandas.DataFrame
        Preprocessed DataFrame containing timeline and heat map entries.
    title : str
        Title to display on the chart.
    zoom_factor : float, optional
        Factor controlling horizontal scaling.  Values <1 compress the
        timeline; values >1 expand it.
    marker_size : int, optional
        Size (area) of milestone diamonds and spot markers.
    suppress_duplicate_labels : bool, optional
        If True, repeated task titles on the same row are drawn only once.

    Returns
    -------
    matplotlib.figure.Figure
        The constructed figure ready for embedding in Tkinter.
    """
    # Separate timeline and heat map
    timeline_df = df[df['Timline / Heat Map'].str.contains('timeline', case=False, na=False)].copy()
    heat_df = df[df['Timline / Heat Map'].str.contains('heat', case=False, na=False)].copy()
    # Row order
    timeline_rows: list[str] = []
    for lr in timeline_df['Line Ref']:
        if pd.isna(lr):
            continue
        if lr not in timeline_rows:
            timeline_rows.append(lr)
    heat_rows: list[str] = []
    for lr in heat_df['Line Ref']:
        if pd.isna(lr):
            continue
        if lr not in heat_rows:
            heat_rows.append(lr)
    # Colour maps
    unique_symbols = list(df['Symbol'].dropna().unique())
    bar_color_map, marker_map, heat_colour_map = generate_color_maps(unique_symbols, heat_rows)
    # Date range
    start_date, end_date = compute_date_range(timeline_df, heat_df)
    start_num = mdates.date2num(start_date)
    end_num = mdates.date2num(end_date)
    total_months = max(1, (end_date.year - start_date.year) * 12 + (end_date.month - start_date.month) + 1)
    width = max(12.0, 0.6 * total_months) * max(0.1, zoom_factor)
    height = max(4.0, 0.6 * max(1, len(timeline_rows)) + 0.4 * max(1, len(heat_rows)) + 1.5)
    fig, (ax_timeline, ax_heat) = plt.subplots(2, 1, figsize=(width, height), sharex=True,
                                               gridspec_kw={'height_ratios': [max(1, len(timeline_rows)), max(1, len(heat_rows))]},
                                               constrained_layout=True)
    bar_height = 0.8
    # Generate a sequence of vertical offsets to stagger labels on the same
    # row.  The scale parameter controls the spacing between labels.  A
    # larger scale increases the vertical separation of overlapping text.
    scale = max(0.0, label_offset_scale)
    label_offsets = [0.0, scale, -scale, 2*scale, -2*scale, 3*scale, -3*scale]
    seen_titles_per_row: dict[str, set] = {lr: set() for lr in timeline_rows}
    for i, line_ref in enumerate(timeline_rows):
        line_tasks = timeline_df[timeline_df['Line Ref'] == line_ref]
        for j, (_, row) in enumerate(line_tasks.iterrows()):
            symbol = row.get('Symbol', '')
            start = row.get('Date From')
            end = row.get('Date To')
            title_txt = str(row.get('Title', '')) if not pd.isna(row.get('Title')) else ''
            if pd.isna(start):
                continue
            start_num_task = mdates.date2num(start)
            if pd.isna(end) or end == start:
                width_days = 1
                end_num_task = start_num_task + width_days
            else:
                end_num_task = mdates.date2num(end)
                width_days = end_num_task - start_num_task
            offset = label_offsets[j % len(label_offsets)]
            # Retrieve per‑task offsets from the DataFrame if present.  These
            # allow individual labels to be shifted horizontally (in days)
            # and vertically (in row units).  Missing values default to 0.
            x_off_val = 0.0
            y_off_val = 0.0
            if 'X Offset' in df.columns:
                try:
                    x_off_val = float(row.get('X Offset', 0) or 0)
                except Exception:
                    x_off_val = 0.0
            if 'Y Offset' in df.columns:
                try:
                    y_off_val = float(row.get('Y Offset', 0) or 0)
                except Exception:
                    y_off_val = 0.0
            # Base y position plus custom per‑task vertical offset
            y_pos = i + offset + y_off_val
            # If the symbol is a custom colour code, treat it as a bar
            if isinstance(symbol, str) and symbol.startswith('#') and len(symbol) in (7, 9):
                # Ensure the custom colour appears in bar_color_map
                bar_color_map.setdefault(symbol, symbol)
                ax_timeline.broken_barh([(start_num_task, width_days)], (i - bar_height / 2, bar_height),
                                        facecolors=symbol, edgecolors='black', linewidth=0.8)
                draw_label = True
                if suppress_duplicate_labels:
                    if title_txt in seen_titles_per_row[line_ref]:
                        draw_label = False
                    else:
                        seen_titles_per_row[line_ref].add(title_txt)
                if draw_label:
                    if width_days >= 3:
                        ax_timeline.text(start_num_task + width_days / 2 + x_off_val,
                                         y_pos,
                                         title_txt,
                                         va='center', ha='center', fontsize=8, color='black', clip_on=True)
                    else:
                        ax_timeline.text(end_num_task + 0.5 + x_off_val,
                                         y_pos,
                                         title_txt,
                                         va='center', ha='left', fontsize=8, color='black', clip_on=True)
            elif symbol in bar_color_map:
                ax_timeline.broken_barh([(start_num_task, width_days)], (i - bar_height / 2, bar_height),
                                        facecolors=bar_color_map[symbol], edgecolors='black', linewidth=0.8)
                draw_label = True
                if suppress_duplicate_labels:
                    if title_txt in seen_titles_per_row[line_ref]:
                        draw_label = False
                    else:
                        seen_titles_per_row[line_ref].add(title_txt)
                if draw_label:
                    if width_days >= 3:
                        ax_timeline.text(start_num_task + width_days / 2 + x_off_val,
                                         y_pos,
                                         title_txt,
                                         va='center', ha='center', fontsize=8, color='black', clip_on=True)
                    else:
                        ax_timeline.text(end_num_task + 0.5 + x_off_val,
                                         y_pos,
                                         title_txt,
                                         va='center', ha='left', fontsize=8, color='black', clip_on=True)
            elif symbol in marker_map:
                specs = marker_map[symbol]
                ax_timeline.scatter(start_num_task, i, marker=specs['marker'], color=specs['color'], s=marker_size, zorder=3)
                if title_txt:
                    draw_label = True
                    if suppress_duplicate_labels:
                        if title_txt in seen_titles_per_row[line_ref]:
                            draw_label = False
                        else:
                            seen_titles_per_row[line_ref].add(title_txt)
                    if draw_label:
                        ax_timeline.text(start_num_task + 1 + x_off_val,
                                         y_pos,
                                         title_txt,
                                         va='center', ha='left', fontsize=8, color='black', clip_on=True)
            else:
                ax_timeline.broken_barh([(start_num_task, width_days)], (i - bar_height / 2, bar_height),
                                        facecolors='#cccccc', edgecolors='black', linewidth=0.8)
                if title_txt:
                    if width_days > 2:
                        ax_timeline.text(start_num_task + width_days / 2 + x_off_val, y_pos, title_txt,
                                         va='center', ha='center', fontsize=8, color='black', clip_on=True)
                    else:
                        ax_timeline.text(end_num_task + 0.5 + x_off_val, y_pos, title_txt,
                                         va='center', ha='left', fontsize=8, color='black', clip_on=True)
    # Configure timeline axes
    ax_timeline.set_xlim(start_num, end_num)
    ax_timeline.set_ylim(-0.5, len(timeline_rows) - 0.5)
    ax_timeline.set_yticks(range(len(timeline_rows)))
    ax_timeline.set_yticklabels(timeline_rows)
    ax_timeline.invert_yaxis()
    ax_timeline.xaxis.set_major_locator(mdates.MonthLocator())
    ax_timeline.xaxis.set_major_formatter(mdates.DateFormatter('%b\n%Y'))
    ax_timeline.tick_params(axis='x', which='major', labelrotation=0)
    ax_timeline.set_title(title, fontsize=14, pad=20)
    ax_timeline.grid(axis='x', which='major', linestyle='--', alpha=0.3)
    # Risk annotations
    for i, line_ref in enumerate(timeline_rows):
        risks = timeline_df[timeline_df['Line Ref'] == line_ref]['Risk Level'].dropna().unique().tolist()
        if not risks:
            continue
        risk_val = risks[0]
        try:
            severity = float(risk_val)
            if severity >= 4.0:
                level_str, colour = 'High', 'red'
            elif severity >= 3.0:
                level_str, colour = 'Medium', '#e6b800'
            else:
                level_str, colour = 'Low', 'green'
        except Exception:
            parts = str(risk_val).split('-')
            level = parts[0].strip().lower()
            if level.startswith('high'):
                level_str, colour = 'High', 'red'
            elif level.startswith('med'):
                level_str, colour = 'Medium', '#e6b800'
            elif level.startswith('low'):
                level_str, colour = 'Low', 'green'
            else:
                level_str, colour = parts[0].strip(), 'grey'
        desc = ''
        if isinstance(risk_val, str) and '-' in risk_val:
            parts_desc = [p.strip() for p in risk_val.split('-', 1)]
            if len(parts_desc) > 1:
                desc = parts_desc[1]
        annotate_text = f"Impact Risk: {level_str}"
        if desc:
            annotate_text += f" - {desc}"
        x_pos = start_num - (end_num - start_num) * 0.02
        ax_timeline.text(x_pos, i - 0.3, annotate_text, va='center', ha='right', fontsize=9, color=colour)
    # Heat map panel
    heat_height = 0.6
    for i, line_ref in enumerate(heat_rows):
        row_events = heat_df[heat_df['Line Ref'] == line_ref]
        for _, row in row_events.iterrows():
            heat_date = row.get('Heat Map Dates')
            if pd.isna(heat_date):
                continue
            month_start = datetime(heat_date.year, heat_date.month, 1)
            month_end = month_start + relativedelta(months=1)
            width_days = (month_end - month_start).days
            start_num_month = mdates.date2num(month_start)
            symbol_heat = str(row.get('Symbol', '')).strip().lower()
            colour_map_simple = {
                'red': '#d62728',
                'black': '#000000',
                'blue': '#1f77b4',
                'grey': '#8c8c8c',
                'gray': '#8c8c8c',
                'white': '#ffffff',
            }
            if symbol_heat in colour_map_simple:
                base_colour = colour_map_simple[symbol_heat]
                facecolor = lighten_color(base_colour, amount=0.3)
            else:
                facecolor = heat_colour_map.get(line_ref, '#dddddd')
            ax_heat.broken_barh([(start_num_month, width_days)], (i - heat_height / 2, heat_height),
                                facecolors=facecolor, edgecolors='none', alpha=0.6)
            ax_heat.text(start_num_month + width_days / 2, i,
                         line_ref, va='center', ha='center', fontsize=7, color='black')
        ax_heat.text(start_num - (end_num - start_num) * 0.01, i,
                     line_ref, va='center', ha='right', fontsize=8)
    ax_heat.set_xlim(start_num, end_num)
    ax_heat.set_ylim(-0.5, len(heat_rows) - 0.5)
    ax_heat.set_yticks([])
    ax_heat.xaxis.set_major_locator(mdates.MonthLocator())
    ax_heat.xaxis.set_major_formatter(mdates.DateFormatter('%b\n%Y'))
    ax_heat.tick_params(axis='x', which='major', labelrotation=0)
    ax_heat.grid(axis='x', which='major', linestyle='--', alpha=0.3)
    for ax in (ax_timeline, ax_heat):
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_visible(False)
    # Legend
    legend_handles: list = []
    used_rows = timeline_df.dropna(subset=['Symbol', 'Title'])
    seen_keys = set()
    for _, row in used_rows.iterrows():
        sym = row['Symbol']
        title_val = str(row['Title'])
        key = (sym, title_val)
        if key in seen_keys:
            continue
        seen_keys.add(key)
        if sym in bar_color_map:
            colour = bar_color_map[sym]
            legend_handles.append(Patch(facecolor=colour, edgecolor='black', label=title_val))
        elif sym in marker_map:
            specs = marker_map[sym]
            legend_handles.append(Line2D([0], [0], marker=specs['marker'], color='w', markerfacecolor=specs['color'],
                                         markersize=8, label=title_val))
        else:
            legend_handles.append(Patch(facecolor='#cccccc', edgecolor='black', label=title_val))
    if legend_handles:
        ncol = min(3, len(legend_handles))
        ax_timeline.legend(handles=legend_handles, loc='upper right', bbox_to_anchor=(1.0, 1.15), ncol=ncol, fontsize=8, frameon=False)
    return fig

# -----------------------------------------------------------------------------
# Tkinter GUI class
#

class GanttGUI:
    """A simple Tkinter application for creating and editing Gantt charts."""

    def __init__(self, root: tk.Tk) -> None:
        self.root = root
        self.root.title("Gantt Chart Creator")
        self.df: pd.DataFrame | None = None
        self.zoom_factor: tk.DoubleVar = tk.DoubleVar(value=1.0)
        self.marker_size: tk.IntVar = tk.IntVar(value=70)
        self.label_offset_scale: tk.DoubleVar = tk.DoubleVar(value=0.25)
        # Dictionaries to hold per‑task offsets in days (x) and rows (y). These
        # allow individual label adjustments without affecting other tasks.  The
        # DataFrame will store 'X Offset' and 'Y Offset' columns to persist
        # these values across edits.
        # Build the interface
        self.build_ui()

    def build_ui(self) -> None:
        # Top frame for file operations and settings
        top_frame = ttk.Frame(self.root)
        top_frame.pack(fill='x', padx=5, pady=5)
        load_btn = ttk.Button(top_frame, text="Load Excel", command=self.load_excel)
        load_btn.pack(side='left', padx=5)
        save_png_btn = ttk.Button(top_frame, text="Save PNG", command=self.save_png)
        save_png_btn.pack(side='left', padx=5)
        export_ppt_btn = ttk.Button(top_frame, text="Export PPT", command=self.export_ppt)
        export_ppt_btn.pack(side='left', padx=5)
        # Zoom slider
        ttk.Label(top_frame, text="Zoom:").pack(side='left', padx=5)
        zoom_scale = ttk.Scale(top_frame, from_=0.5, to=2.0, orient='horizontal', variable=self.zoom_factor, command=self.on_zoom_change)
        zoom_scale.pack(side='left', padx=5, fill='x', expand=True)
        ttk.Label(top_frame, text="Marker size:").pack(side='left', padx=5)
        marker_scale = ttk.Scale(top_frame, from_=50, to=150, orient='horizontal', variable=self.marker_size, command=self.on_marker_change)
        marker_scale.pack(side='left', padx=5)
        # Label offset slider
        ttk.Label(top_frame, text="Label offset:").pack(side='left', padx=5)
        offset_scale = ttk.Scale(top_frame, from_=0.0, to=1.0, orient='horizontal', variable=self.label_offset_scale, command=self.on_offset_change)
        offset_scale.pack(side='left', padx=5, fill='x', expand=False)
        # Middle frame: Treeview for task editing
        mid_frame = ttk.Frame(self.root)
        mid_frame.pack(fill='both', expand=True, padx=5, pady=5)
        # Define more columns to display key fields for editing
        columns = ('index', 'Timline / Heat Map', 'Line Ref', 'Title', 'Symbol', 'Risk Level', 'Date From', 'Date To', 'Heat Map Dates')
        self.tree = ttk.Treeview(mid_frame, columns=columns, show='headings', selectmode='browse')
        for col in columns:
            self.tree.heading(col, text=col)
            # Set narrower width for date columns
            if col == 'index':
                width = 60
            elif col in ('Date From', 'Date To', 'Heat Map Dates'):
                width = 90
            else:
                width = 120
            self.tree.column(col, width=width, anchor='w')
        self.tree.pack(side='left', fill='both', expand=True)
        # Scrollbar
        scrollbar = ttk.Scrollbar(mid_frame, orient='vertical', command=self.tree.yview)
        scrollbar.pack(side='right', fill='y')
        self.tree.configure(yscrollcommand=scrollbar.set)
        # Editing panel with more fields
        edit_frame = ttk.LabelFrame(self.root, text="Edit Task")
        edit_frame.pack(fill='x', padx=5, pady=5)
        row_idx = 0
        ttk.Label(edit_frame, text="Title:").grid(row=row_idx, column=0, sticky='w')
        self.edit_title_var = tk.StringVar()
        self.edit_title_entry = ttk.Entry(edit_frame, textvariable=self.edit_title_var)
        self.edit_title_entry.grid(row=row_idx, column=1, sticky='ew', padx=5)
        row_idx += 1
        ttk.Label(edit_frame, text="Symbol:").grid(row=row_idx, column=0, sticky='w')
        self.edit_symbol_var = tk.StringVar()
        self.edit_symbol_combo = ttk.Combobox(edit_frame, textvariable=self.edit_symbol_var, state='readonly')
        # Populate later when data is loaded
        self.edit_symbol_combo.grid(row=row_idx, column=1, sticky='ew', padx=5)
        row_idx += 1
        ttk.Label(edit_frame, text="Custom Colour:").grid(row=row_idx, column=0, sticky='w')
        self.edit_colour_var = tk.StringVar()
        self.edit_colour_btn = ttk.Button(edit_frame, text="Choose", command=self.choose_colour)
        self.edit_colour_btn.grid(row=row_idx, column=1, sticky='w', padx=5)
        row_idx += 1
        ttk.Label(edit_frame, text="Risk Level:").grid(row=row_idx, column=0, sticky='w')
        self.edit_risk_var = tk.StringVar()
        self.edit_risk_entry = ttk.Entry(edit_frame, textvariable=self.edit_risk_var)
        self.edit_risk_entry.grid(row=row_idx, column=1, sticky='ew', padx=5)
        row_idx += 1
        ttk.Label(edit_frame, text="Date From (YYYY‑MM‑DD):").grid(row=row_idx, column=0, sticky='w')
        self.edit_date_from_var = tk.StringVar()
        self.edit_date_from_entry = ttk.Entry(edit_frame, textvariable=self.edit_date_from_var)
        self.edit_date_from_entry.grid(row=row_idx, column=1, sticky='ew', padx=5)
        row_idx += 1
        ttk.Label(edit_frame, text="Date To (YYYY‑MM‑DD):").grid(row=row_idx, column=0, sticky='w')
        self.edit_date_to_var = tk.StringVar()
        self.edit_date_to_entry = ttk.Entry(edit_frame, textvariable=self.edit_date_to_var)
        self.edit_date_to_entry.grid(row=row_idx, column=1, sticky='ew', padx=5)
        row_idx += 1
        ttk.Label(edit_frame, text="Heat Map Date (YYYY‑MM‑DD):").grid(row=row_idx, column=0, sticky='w')
        self.edit_heat_date_var = tk.StringVar()
        self.edit_heat_date_entry = ttk.Entry(edit_frame, textvariable=self.edit_heat_date_var)
        self.edit_heat_date_entry.grid(row=row_idx, column=1, sticky='ew', padx=5)
        row_idx += 1
        # X and Y label offsets allow individual adjustment of the label
        # position relative to its default.  X Offset is in days (positive
        # values move the label to the right), and Y Offset is in units of
        # row spacing (positive values move the label downward on an inverted
        # y‑axis).
        ttk.Label(edit_frame, text="X Offset (days):").grid(row=row_idx, column=0, sticky='w')
        self.edit_x_offset_var = tk.StringVar()
        self.edit_x_offset_entry = ttk.Entry(edit_frame, textvariable=self.edit_x_offset_var)
        self.edit_x_offset_entry.grid(row=row_idx, column=1, sticky='ew', padx=5)
        row_idx += 1
        ttk.Label(edit_frame, text="Y Offset (rows):").grid(row=row_idx, column=0, sticky='w')
        self.edit_y_offset_var = tk.StringVar()
        self.edit_y_offset_entry = ttk.Entry(edit_frame, textvariable=self.edit_y_offset_var)
        self.edit_y_offset_entry.grid(row=row_idx, column=1, sticky='ew', padx=5)
        row_idx += 1
        update_btn = ttk.Button(edit_frame, text="Update Task", command=self.update_task)
        update_btn.grid(row=row_idx, column=0, columnspan=2, pady=5)
        # Chart area
        self.fig = None
        self.canvas = None
        self.chart_frame = ttk.Frame(self.root)
        self.chart_frame.pack(fill='both', expand=True, padx=5, pady=5)
        # Bind tree selection
        self.tree.bind('<<TreeviewSelect>>', self.on_tree_select)

    # -------------------------------------------------------------------------
    # File operations
    def load_excel(self) -> None:
        path = filedialog.askopenfilename(title="Select Excel file", filetypes=[("Excel files", "*.xlsx;*.xls")])
        if not path:
            return
        try:
            df_raw = pd.read_excel(path, sheet_name=0)
            df_raw.columns = df_raw.columns.str.strip()
            # Parse dates
            date_cols = ['Date From', 'Date To', 'Heat Map Dates']
            for col in date_cols:
                if col in df_raw.columns:
                    df_raw[col] = pd.to_datetime(df_raw[col], errors='coerce')
            self.df = df_raw
            self.populate_tree()
            # Update symbol combobox with available symbols plus 'Custom'
            unique_symbols = sorted({str(s) for s in self.df['Symbol'].dropna().unique()})
            options = [s for s in unique_symbols]
            if 'Custom' not in options:
                options.append('Custom')
            self.edit_symbol_combo['values'] = options
            self.render_chart()
        except Exception as e:
            messagebox.showerror("Error", f"Failed to load Excel file:\n{e}")

    def save_png(self) -> None:
        if self.fig is None:
            messagebox.showwarning("No Chart", "Please generate a chart first.")
            return
        path = filedialog.asksaveasfilename(title="Save PNG", defaultextension=".png", filetypes=[("PNG Files", "*.png")])
        if not path:
            return
        try:
            self.fig.savefig(path, dpi=300, bbox_inches='tight')
            messagebox.showinfo("Saved", f"Chart saved to {path}")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to save PNG:\n{e}")

    def export_ppt(self) -> None:
        if self.fig is None:
            messagebox.showwarning("No Chart", "Please generate a chart first.")
            return
        # Save temporary image
        tmp_path = os.path.join(os.getcwd(), 'tmp_gantt.png')
        self.fig.savefig(tmp_path, dpi=300, bbox_inches='tight')
        ppt_path = filedialog.asksaveasfilename(title="Save PowerPoint", defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
        if not ppt_path:
            os.remove(tmp_path)
            return
        try:
            self._export_to_powerpoint(tmp_path, ppt_path)
            os.remove(tmp_path)
            messagebox.showinfo("Saved", f"PowerPoint saved to {ppt_path}")
        except Exception as e:
            os.remove(tmp_path)
            messagebox.showerror("Error", f"Failed to export to PowerPoint:\n{e}")

    def _export_to_powerpoint(self, image_path: str, ppt_path: str) -> None:
        """Embed a saved image into a PowerPoint file.  This method mirrors
        the export_to_powerpoint function from gantt_creator.py but is
        self‑contained so that the GUI can export without external imports.
        """
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        im = Image.open(image_path)
        width, height = im.size
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        margin = Inches(0.3)
        available_width = slide_width - 2 * margin
        available_height = slide_height - 2 * margin
        img_ratio = width / height
        slide_ratio = available_width / available_height
        if img_ratio > slide_ratio:
            pic_width = available_width
            pic_height = available_width / img_ratio
        else:
            pic_height = available_height
            pic_width = available_height * img_ratio
        left = margin + (available_width - pic_width) / 2
        top = margin + (available_height - pic_height) / 2
        slide.shapes.add_picture(image_path, left, top, width=pic_width, height=pic_height)
        prs.save(ppt_path)

    # -------------------------------------------------------------------------
    # Data and chart operations
    def populate_tree(self) -> None:
        """Populate the Treeview with rows from the DataFrame."""
        for row in self.tree.get_children():
            self.tree.delete(row)
        if self.df is None:
            return
        for idx, row in self.df.iterrows():
            # Format dates for display
            def fmt(val):
                return '' if pd.isna(val) else val.strftime('%Y-%m-%d')
            values = (
                idx,
                row.get('Timline / Heat Map', ''),
                row.get('Line Ref', ''),
                row.get('Title', ''),
                row.get('Symbol', ''),
                row.get('Risk Level', ''),
                fmt(row.get('Date From')),
                fmt(row.get('Date To')),
                fmt(row.get('Heat Map Dates')),
            )
            self.tree.insert('', 'end', iid=str(idx), values=values)
    def on_tree_select(self, event: tk.Event) -> None:
        selected = self.tree.selection()
        if not selected or self.df is None:
            return
        idx = int(selected[0])
        row = self.df.loc[idx]
        # Populate form fields
        self.edit_title_var.set(str(row.get('Title', '')))
        sym_val = str(row.get('Symbol', ''))
        self.edit_symbol_var.set(sym_val)
        # If symbol looks like a hex colour, store in colour var
        if sym_val.startswith('#'):
            self.edit_colour_var.set(sym_val)
        else:
            self.edit_colour_var.set('')
        self.edit_risk_var.set(str(row.get('Risk Level', '')))
        # Dates
        def fmt_date(val):
            return '' if pd.isna(val) else val.strftime('%Y-%m-%d')
        self.edit_date_from_var.set(fmt_date(row.get('Date From')))
        self.edit_date_to_var.set(fmt_date(row.get('Date To')))
        self.edit_heat_date_var.set(fmt_date(row.get('Heat Map Dates')))
        # X and Y offsets
        def fmt_offset(val):
            try:
                return '' if pd.isna(val) else str(float(val))
            except Exception:
                return ''
        # Ensure offset columns exist in DataFrame; if missing, default to 0
        if 'X Offset' not in self.df.columns:
            self.df['X Offset'] = 0.0
        if 'Y Offset' not in self.df.columns:
            self.df['Y Offset'] = 0.0
        self.edit_x_offset_var.set(fmt_offset(row.get('X Offset')))
        self.edit_y_offset_var.set(fmt_offset(row.get('Y Offset')))
    def choose_colour(self) -> None:
        """Open a colour chooser and store the result in edit_colour_var."""
        colour_code = colorchooser.askcolor(title="Choose bar colour")
        if colour_code and colour_code[1]:
            # colour_code[1] is the hex string (e.g. '#f05ded').  When a custom
            # colour is chosen, automatically set the symbol variable to the
            # hex string so the update routine recognises it as a bar colour.
            hex_str = colour_code[1]
            self.edit_colour_var.set(hex_str)
            # Set symbol combobox to the hex colour (treated as custom symbol)
            self.edit_symbol_var.set(hex_str)
    def update_task(self) -> None:
        """Apply edits from the form to the selected task."""
        selected = self.tree.selection()
        if not selected or self.df is None:
            messagebox.showwarning("No Selection", "Please select a task to update.")
            return
        idx = int(selected[0])
        # Update fields in DataFrame
        # Title
        self.df.at[idx, 'Title'] = self.edit_title_var.get()
        # Symbol handling: either a known symbol or custom colour
        chosen_symbol = self.edit_symbol_var.get()
        colour_code = self.edit_colour_var.get()
        if chosen_symbol.lower() == 'custom' or (chosen_symbol.startswith('#') and len(chosen_symbol) in (7, 9)):
            # Use the colour code if present; fallback to the chosen symbol if it's a hex
            if colour_code.startswith('#'):
                self.df.at[idx, 'Symbol'] = colour_code
            elif chosen_symbol.startswith('#'):
                self.df.at[idx, 'Symbol'] = chosen_symbol
            else:
                messagebox.showwarning("Custom Colour", "Please choose a custom colour using the Colour button.")
        else:
            self.df.at[idx, 'Symbol'] = chosen_symbol
        # Risk level
        self.df.at[idx, 'Risk Level'] = self.edit_risk_var.get()
        # Dates
        def parse_date(s: str):
            s = s.strip()
            return pd.to_datetime(s) if s else pd.NaT
        self.df.at[idx, 'Date From'] = parse_date(self.edit_date_from_var.get())
        self.df.at[idx, 'Date To'] = parse_date(self.edit_date_to_var.get())
        self.df.at[idx, 'Heat Map Dates'] = parse_date(self.edit_heat_date_var.get())
        # Update per‑task label offsets.  Parse numeric values; if blank or invalid
        # revert to 0.0.
        try:
            x_off = float(self.edit_x_offset_var.get().strip()) if self.edit_x_offset_var.get().strip() else 0.0
        except ValueError:
            x_off = 0.0
        try:
            y_off = float(self.edit_y_offset_var.get().strip()) if self.edit_y_offset_var.get().strip() else 0.0
        except ValueError:
            y_off = 0.0
        # Ensure offset columns exist
        if 'X Offset' not in self.df.columns:
            self.df['X Offset'] = 0.0
        if 'Y Offset' not in self.df.columns:
            self.df['Y Offset'] = 0.0
        self.df.at[idx, 'X Offset'] = x_off
        self.df.at[idx, 'Y Offset'] = y_off
        # Update treeview row with new values
        row_vals = (
            idx,
            self.df.at[idx, 'Timline / Heat Map'],
            self.df.at[idx, 'Line Ref'],
            self.df.at[idx, 'Title'],
            self.df.at[idx, 'Symbol'],
            self.df.at[idx, 'Risk Level'],
            self.df.at[idx, 'Date From'].strftime('%Y-%m-%d') if not pd.isna(self.df.at[idx, 'Date From']) else '',
            self.df.at[idx, 'Date To'].strftime('%Y-%m-%d') if not pd.isna(self.df.at[idx, 'Date To']) else '',
            self.df.at[idx, 'Heat Map Dates'].strftime('%Y-%m-%d') if not pd.isna(self.df.at[idx, 'Heat Map Dates']) else ''
        )
        # Optionally update tree with offset values? Offsets are not displayed in tree
        self.tree.item(str(idx), values=row_vals)
        # Re-render chart
        self.render_chart()
    def on_zoom_change(self, event: str | None = None) -> None:
        # event is a string representing the slider value
        try:
            if event is not None:
                self.zoom_factor.set(float(event))
        except ValueError:
            pass
        if self.df is not None:
            self.render_chart()
    def on_marker_change(self, event: str | None = None) -> None:
        try:
            if event is not None:
                self.marker_size.set(int(float(event)))
        except ValueError:
            pass
        if self.df is not None:
            self.render_chart()

    def on_offset_change(self, event: str | None = None) -> None:
        """Callback for label offset slider.  Update internal variable and redraw."""
        try:
            if event is not None:
                self.label_offset_scale.set(float(event))
        except ValueError:
            pass
        if self.df is not None:
            self.render_chart()
    def render_chart(self) -> None:
        """Regenerate the chart using current settings and display it."""
        if self.df is None:
            return
        # Create figure
        title = "Gantt Chart"
        fig = create_gantt_figure(
            self.df,
            title=title,
            zoom_factor=self.zoom_factor.get(),
            marker_size=self.marker_size.get(),
            label_offset_scale=self.label_offset_scale.get(),
        )
        # Remove old canvas
        if self.canvas is not None:
            self.canvas.get_tk_widget().destroy()
        self.fig = fig
        self.canvas = FigureCanvasTkAgg(fig, master=self.chart_frame)
        self.canvas.draw()
        self.canvas.get_tk_widget().pack(fill='both', expand=True)


def main() -> None:
    root = tk.Tk()
    app = GanttGUI(root)
    root.mainloop()


if __name__ == '__main__':
    main()